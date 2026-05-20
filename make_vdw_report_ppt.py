from pathlib import Path
import csv

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
DESKTOP = Path(r"C:\Users\25739\Desktop")
OUT = DESKTOP / "EFT-vdW汇报.pptx"
IMG_DIR = ROOT / "_ppt_assets"
IMG_DIR.mkdir(exist_ok=True)

FONT = "Microsoft YaHei"
BLUE = RGBColor(31, 78, 121)
DARK = RGBColor(35, 35, 35)
LIGHT = RGBColor(245, 248, 252)
ACCENT = RGBColor(192, 80, 77)
GREEN = RGBColor(79, 129, 89)


def read_csv(path):
    with open(path, newline="", encoding="utf-8") as f:
        return list(csv.DictReader(f))


tdhf_nstates = read_csv(ROOT / "results/ar/ar_tdhf_nstates_convergence.csv")
mo_basis = read_csv(ROOT / "results/ar/ar_mo_basis_convergence.csv")
tdhf_basis = read_csv(ROOT / "results/ar/ar_tdhf_basis_convergence.csv")
tail_rows = read_csv(ROOT / "results/ar/ar2_tail_comparison.csv")


def prs_init():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def set_bg(slide, color=RGBColor(255, 255, 255)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, x, y, w, h, text="", font_size=24, bold=False, color=DARK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.55, 0.32, 12.2, 0.55, title, 25, True, BLUE)
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(0.93), Inches(12.2), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.color.rgb = BLUE
    if subtitle:
        add_textbox(slide, 0.65, 1.03, 11.8, 0.35, subtitle, 12, False, RGBColor(100, 100, 100))


def add_bullets(slide, x, y, w, h, bullets, font_size=18, color=DARK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = item
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return box


def add_code(slide, x, y, w, h, text, font_size=15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 242, 245)
    box.line.color.rgb = RGBColor(220, 220, 220)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Consolas"
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor(40, 40, 40)
    return box


def add_table(slide, x, y, w, h, headers, rows, font_size=12):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for col, head in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = FONT
                r.font.size = Pt(font_size)
                r.font.bold = True
                r.font.color.rgb = RGBColor(255, 255, 255)
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = Pt(font_size)
                    r.font.color.rgb = DARK
    return table


def save_tdhf_nstates_chart():
    fig, ax = plt.subplots(figsize=(7.4, 4.1), dpi=160)
    for basis, marker in [("aug-cc-pVTZ", "o"), ("aug-cc-pVQZ", "s")]:
        rows = [r for r in tdhf_nstates if r["basis"] == basis]
        x = [int(r["nstates"]) for r in rows]
        c6 = [float(r["C6"]) for r in rows]
        ax.plot(x, c6, marker=marker, label=basis)
    ax.axhline(64.3, color="gray", linestyle="--", linewidth=1, label="reference C6=64.3")
    ax.set_xlabel("TDHF nstates")
    ax.set_ylabel("Ar-Ar C6 (a.u.)")
    ax.set_title("TDHF C6 convergence")
    ax.legend()
    ax.grid(alpha=0.25)
    path = IMG_DIR / "tdhf_nstates_c6.png"
    fig.tight_layout()
    fig.savefig(path)
    plt.close(fig)
    return path


def save_method_c6_chart():
    labels = ["Reference", "Calibrated", "MO\naug-QZ", "TDHF\naug-QZ\n200"]
    values = [64.3, 64.3, 76.13748508, 60.73027908]
    colors = ["#808080", "#4F81BD", "#C0504D", "#4F8159"]
    fig, ax = plt.subplots(figsize=(7.2, 4.0), dpi=160)
    bars = ax.bar(labels, values, color=colors)
    ax.axhline(64.3, color="black", linewidth=0.8, linestyle="--")
    ax.set_ylabel("C6 (a.u.)")
    ax.set_title("Ar C6 method comparison")
    ax.set_ylim(0, 85)
    for bar, value in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, value + 1, f"{value:.2f}", ha="center", fontsize=9)
    ax.grid(axis="y", alpha=0.2)
    path = IMG_DIR / "method_c6.png"
    fig.tight_layout()
    fig.savefig(path)
    plt.close(fig)
    return path


def save_tail_chart():
    fig, ax = plt.subplots(figsize=(7.4, 4.1), dpi=160)
    r = [float(x["R_bohr"]) for x in tail_rows]
    ax.plot(r, [abs(float(x["E_ref_meV"])) for x in tail_rows], "k--", marker="o", label="Reference")
    ax.plot(r, [abs(float(x["E_tdhf_meV"])) for x in tail_rows], marker="s", label="TDHF")
    ax.plot(r, [abs(float(x["E_mo_meV"])) for x in tail_rows], marker="^", label="MO")
    ax.set_yscale("log")
    ax.set_xlabel("R (Bohr)")
    ax.set_ylabel("|E| (meV), log scale")
    ax.set_title("Ar2 long-range tail: |E| = C6/R^6")
    ax.legend()
    ax.grid(alpha=0.25, which="both")
    path = IMG_DIR / "ar2_tail.png"
    fig.tight_layout()
    fig.savefig(path)
    plt.close(fig)
    return path


def add_picture(slide, path, x, y, w, h):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def make_ppt():
    prs = prs_init()
    blank = prs.slide_layouts[6]

    # 1 Title
    slide = prs.slides.add_slide(blank)
    set_bg(slide, RGBColor(248, 250, 253))
    add_textbox(slide, 0.7, 0.7, 12.0, 0.8, "EFT-vdW：从 frozen-core 动力学到 Ar 的 TDHF C6 benchmark", 28, True, BLUE)
    add_textbox(slide, 0.75, 1.65, 11.5, 0.55, "基于《算vdW》推导、PRL EFT-KS 论文和当前代码结果", 18, False, DARK)
    add_textbox(slide, 0.75, 6.55, 11.5, 0.35, "当前里程碑：TDHF aug-cc-pVQZ nstates=200，Ar C6 error = -5.55%", 15, True, GREEN)

    # 2 EFT-vdW objective
    slide = prs.slides.add_slide(blank)
    add_title(slide, "1. EFT-vdW 要解决什么问题？", "把 frozen-core 动力学从 band narrowing 推广到长程 correlation")
    add_bullets(slide, 0.7, 1.25, 6.0, 5.2, [
        "传统赝势把 core 当成静态对象：只留下静态 Hartree/Fock/投影效应。",
        "PRL 论文的关键发现：积分掉 core 后还剩一个动态 pole， conventional pseudopotential 没有。",
        "vdW 本质是电荷涨落的长程相关；因此同一个 core 动态 pole 应该也贡献 core polarizability。",
        "EFT-vdW 的目标：从可匹配的 core Wilson coefficient 出发，构造 α(iξ) 和 C6，而不是只拟合经验参数。"
    ], 17)
    add_code(slide, 7.0, 1.35, 5.6, 4.3,
             "PRL:  core dynamics → z_core → band narrowing\n\n"
             "本工作: core dynamics → χ_c(iξ) → α(iξ) → C6\n\n"
             "核心问题：\n"
             "static pseudopotential 已经包含什么？\n"
             "dynamic fluctuation 还剩什么？\n"
             "如何避免 double counting？", 16)

    # 3 Source-functional derivation
    slide = prs.slides.add_slide(blank)
    add_title(slide, "2. 第一步：带外源积分掉 core")
    add_bullets(slide, 0.65, 1.2, 5.85, 5.4, [
        "给 core density 加标量外源 Φ，Φ 包含外场和 valence density 产生的 Coulomb 势。",
        "对 core 自由度做 cumulant expansion：线性项/静态项进入普通赝势；二次项就是动态 core response。",
        "PRL Eq.(6) 是这个二次核投影到 valence one-particle self-energy 后的一粒子版本。",
        "因此 vdW 不需要重新发明新对象，而是把同一个 pole 解释成 density response 的谱分解。"
    ], 16)
    add_code(slide, 6.75, 1.1, 5.95, 5.05,
             "Φ_A(r,τ) = φ_ext(r,τ) + ∫dr' v(r-r') ρ_v(r',τ)\n\n"
             "S_eff_core[Φ] = S0 + ∫ n_c^0 Φ\n"
             "                 - 1/2 ∫ Φ χ_c Φ + ...\n\n"
             "静态部分 → V_PSP / V_xc matching\n"
             "动态二次核 → core fluctuation / vdW", 14)

    # 4 Spectral pole and form factor
    slide = prs.slides.add_slide(blank)
    add_title(slide, "3. 第二步：把 Eq.(6)/(7) 写成 core response 的谱表示")
    add_bullets(slide, 0.65, 1.25, 5.85, 5.15, [
        "PRL Eq.(6)：动态项有 pole 结构 1/(iω+ΔE_c)。",
        "响应函数需要偶的 Matsubara 组合：1/(iξ+Δ)+1/(-iξ+Δ)=2Δ/(ξ²+Δ²)。",
        "Eq.(7) 的 Coulomb-dressed form factor f_K^c 可还原 transition density τ(q)。",
        "vdW leading term 是偶极-偶极涨落，所以需要从 scalar channel 推广到 ℓ=1 dipole channel。"
    ], 16)
    add_code(slide, 6.7, 1.0, 6.0, 5.35,
             "τ_λ(q) = q² f_λ(q) / 4π\n\n"
             "χ_c(q,q'; iξ) = Σ_λ [2Δ_λ τ_λ(q)τ*_λ(q')]\n"
             "                 / [ξ² + Δ_λ²]\n\n"
             "vχ_cv = Σ_λ [2Δ_λ f_λ(q)f*_λ(q')]\n"
             "       / [ξ² + Δ_λ²]\n\n"
             "这一步把 one-particle dynamic PSP\n"
             "提升为 density response kernel。", 14)

    # 5 Dipole limit
    slide = prs.slides.add_slide(blank)
    add_title(slide, "4. 第三步：长波偶极极限得到 α(iξ)")
    add_bullets(slide, 0.65, 1.25, 5.85, 5.1, [
        "长程 vdW 只需要 q→0 的 multipole 展开；leading term 是 dipole。",
        "transition dipole 由 transition density 的一阶导数给出。",
        "闭壳层各向同性时，张量极化率化成 scalar α(iξ)。",
        "代码里统一用 oscillator strength 形式：α(iξ)=Σ osc/(Δ²+ξ²)。"
    ], 16)
    add_code(slide, 6.75, 1.15, 5.95, 4.9,
             "d_{λ,i} = i ∂τ_λ(q)/∂q_i |_{q=0}\n\n"
             "α_{ij}(iξ) = Σ_λ 2Δ_λ d_{λ,i}d*_{λ,j}\n"
             "             / (Δ_λ² + ξ²)\n\n"
             "α(iξ) = (1/3)Tr α_ij\n"
             "      = Σ_λ osc_λ/(Δ_λ²+ξ²)\n\n"
             "osc_λ = (2/3) Δ_λ |d_λ|²", 15)

    # 6 Energy and screening
    slide = prs.slides.add_slide(blank)
    add_title(slide, "5. 第四步：从 α(iξ) 到 C6 与 screened vdW energy")
    add_bullets(slide, 0.65, 1.25, 5.9, 5.25, [
        "真空中两原子远距离极限恢复 London 公式：E=-C6/R⁶。",
        "固体中不能只用裸 Coulomb，需要 valence-screened interaction W_v。",
        "EFT 给出的是统一的 RPA/log correlation 结构；pairwise C6 是二阶、远距、真空极限。",
        "因此当前 Ar benchmark 是最小闭环：先验证 α(iξ)→C6→tail，再进入屏蔽/多体版本。"
    ], 16)
    add_code(slide, 6.7, 1.0, 6.05, 5.45,
             "C6_AB = (3/π) ∫₀∞ dξ α_A(iξ) α_B(iξ)\n\n"
             "E_AB(R) = - C6_AB / R_AB⁶\n\n"
             "固体/凝聚态推广：\n"
             "W_v(iξ) = [v⁻¹ - χ_v^irr(iξ)]⁻¹\n\n"
             "E_vdW = (1/2π)∫dξ Tr_inter{ ln[1-W_vχ_c]\n"
             "                         + W_vχ_c }", 14)

    # 7 Double counting
    slide = prs.slides.add_slide(blank)
    add_title(slide, "6. Double counting：为什么 EFT 写法更干净")
    add_bullets(slide, 0.65, 1.25, 12.0, 4.8, [
        "静态 core effect 已经在 H_KS、V_xc 和 pseudopotential 中匹配过，不能再加一次。",
        "动态 self-energy 要做 subtraction：只保留相对匹配点的 frequency-dependent residual。",
        "vdW energy 只取 inter-center fluctuation：去掉 A=A self block 和本地 Wilson coefficient。",
        "如果 semicore 已经显式放进 valence，必须从 core response χ_c 中删掉，避免 core/valence 重复。"
    ], 18)
    add_code(slide, 1.0, 5.35, 11.4, 0.8,
             "EFT 原则：tree-level KS 已经包含的 local/relevant/marginal pieces 不再重复；只计算 residual nonlocal dynamical correlation。", 16)

    # 8 Traditional comparison
    slide = prs.slides.add_slide(blank)
    add_title(slide, "7. 与传统 vdW 做法的区别")
    add_table(slide, 0.45, 1.05, 12.45, 4.25,
              ["方法", "输入/参数", "优点", "局限"],
              [
                  ["DFT-D4", "经验原子参数、charge-dependent scaling", "稳健、便宜、工程可用", "不是从 core Wilson coefficient 推导"],
                  ["MBD", "原子极化率、C6、R_vdW + screening", "有 many-body screening", "参数通常是 total atom fitted data"],
                  ["RPA/ACFDT", "KS response + Coulomb kernel", "原则上第一性原理", "贵；core/PSP double counting 不透明"],
                  ["EFT-vdW", "core excitation/form factor → α(iξ)", "可追溯 core dynamics；自然处理 double counting", "需要可靠 oscillator spectrum 与 screening"],
              ], 10)
    add_bullets(slide, 0.8, 5.65, 11.8, 0.8, [
        "核心优势不是“更便宜地拟合 C6”，而是把 pseudopotential 漏掉的 core dynamic response 作为可匹配 Wilson coefficient 接回 vdW。"
    ], 15, GREEN)

    # 9 Advantages
    slide = prs.slides.add_slide(blank)
    add_title(slide, "8. EFT-vdW 的潜在优势")
    add_bullets(slide, 0.8, 1.2, 11.8, 5.4, [
        "1. 可解释性：每个 channel 有 Δ、osc、shell 来源，可做 TRK 和 channel contribution 审计。",
        "2. double-counting 清晰：静态 PSP / V_xc 与动态 residual 分层处理。",
        "3. 可系统改进：radial → 3D MO → TDHF/RPA → screened W_v，而不是一次性经验拟合。",
        "4. 可连接传统方法：EFT-derived α0/C6 可作为 MBD/D4 参数替换或 core correction。",
        "5. 可服务 pseudopotential 问题：回答 valence-only 会漏多少 core response，以及 EFT core channel 补回多少。"
    ], 18)

    # 2 Motivation
    slide = prs.slides.add_slide(blank)
    add_title(slide, "9. 物理动机：PRL 里的 frozen-core 动力学")
    add_bullets(slide, 0.75, 1.35, 5.9, 4.9, [
        "论文核心：conventional pseudopotential 冻结了 core 的动态响应。",
        "Eq.(6) 给出 core 动力学 pole：M_sc M_tc* / (iω + ΔE_c)。",
        "Eq.(7) 给出 Coulomb-dressed core form factor f_K^c。",
        "原文用它解释 KS band narrowing；vdW 需要把同一动力学项解释为 core density response。"
    ], 18)
    add_code(slide, 7.05, 1.45, 5.55, 3.45,
             "V_dyn(K,K'; iω) = Σ_c f_K^c f_K'^c / (iω + ΔE_c)\n\n"
             "关键转译：\ncore dynamical pole\n→ density response χ_c(q,q';iξ)\n→ dipole polarizability α(iξ)\n→ C6")

    # 3 Derivation bridge
    slide = prs.slides.add_slide(blank)
    add_title(slide, "10. 从 form factor 到 vdW：推导桥梁")
    add_bullets(slide, 0.7, 1.15, 5.7, 5.5, [
        "引入外源 φ，core 被积分掉后出现二次响应核 χ_c。",
        "Coulomb-dressed form factor 与 transition density 的关系：τ(q)=q² f(q)/(4π)。",
        "长波极限取 dipole：d_i = i ∂τ(q)/∂q_i |_{q=0}。",
        "闭壳层各向同性原子：α(iξ)=Σλ oscλ/(Δλ²+ξ²)。"
    ], 17)
    add_code(slide, 6.8, 1.15, 5.8, 4.75,
             "χ_c(q,q'; iξ) = Σ_λ 2Δ_λ τ_λ(q)τ_λ*(q') / (ξ²+Δ_λ²)\n\n"
             "osc_λ = (2/3) Δ_λ |d_λ|²\n\n"
             "C6_AB = (3/π) ∫₀∞ α_A(iξ) α_B(iξ) dξ\n\n"
             "E_AB(R) = - C6_AB / R⁶", 14)

    # 4 Implementation chain
    slide = prs.slides.add_slide(blank)
    add_title(slide, "11. 代码实现链条")
    add_code(slide, 0.85, 1.25, 11.8, 4.7,
             "输入层：\n"
             "  radial_orbitals.csv / MO dipoles / TDHF oscillator strengths\n\n"
             "channel 层：\n"
             "  atom, channel, delta_Ha, osc, is_core, source\n\n"
             "后端层：\n"
             "  eft_alpha.py: alpha_iw_from_osc(), c6_from_alpha(), pairwise_energy()\n\n"
             "benchmark 层：\n"
             "  run_alpha_table.py --input ...\n"
             "  run_c6_table.py --input ...\n"
             "  compare_alpha_c6.py\n"
             "  run_ar2_tail.py", 16)

    # 5 Routes
    slide = prs.slides.add_slide(blank)
    add_title(slide, "12. 三条 oscillator 路线：从诊断到 baseline")
    add_table(slide, 0.55, 1.35, 12.25, 3.1,
              ["路线", "输入", "用途", "当前状态"],
              [
                  ["Radial shell-average", "u(r) → radial d²", "诊断 EFT dipole 接口", "存在 3p shell TRK 过计数"],
                  ["3D MO", "AO dipole integral + MO coeff", "当前 prediction baseline", "TRK 不爆，C6 error +18%(MO QZ)"],
                  ["TDHF/RPA", "PySCF TDHF Δλ, fλ", "当前最佳 Ar benchmark", "C6 error -5.55%(QZ, 200 states)"],
              ], 12)
    add_bullets(slide, 0.9, 4.85, 11.8, 1.1, [
        "重要诊断：radial aug-cc-pVQZ 的 sum_osc_discrete/N_core = 1.7118，不能作为 prediction。",
        "3D MO 和 TDHF 直接使用 PySCF 矩阵元/响应，绕开 shell-average 简并过计数。"
    ], 16)

    # 6 Radial audit
    slide = prs.slides.add_slide(blank)
    add_title(slide, "13. Radial 路线审计：问题定位到 shell-average")
    add_bullets(slide, 0.7, 1.25, 5.9, 4.7, [
        "per-shell TRK 检查显示：Ar 3p shell 是主要过计数来源。",
        "aug-cc-pVQZ radial：3p sum_osc / occupation = 4.21。",
        "同 basis 的 3D MO TRK：sum_f/N = 0.9407，没有爆掉。",
        "结论：问题不是 PySCF basis 本身，而是 radial shell-average / shell degeneracy 路径。"
    ], 17)
    add_table(slide, 7.0, 1.45, 5.5, 2.6,
              ["检查", "结果"],
              [
                  ["radial total", "sum_osc/N = 1.7118"],
                  ["radial 3p shell", "ratio = 4.2103"],
                  ["3D MO total", "sum_f/N = 0.9407"],
                  ["判定", "radial route diagnostic only"],
              ], 12)

    # 7 Method C6 chart
    slide = prs.slides.add_slide(blank)
    add_title(slide, "14. Ar C6 方法对比")
    add_picture(slide, save_method_c6_chart(), 0.75, 1.3, 6.0, 3.6)
    add_table(slide, 7.05, 1.25, 5.65, 3.35,
              ["method", "C6", "error", "role"],
              [
                  ["reference", "64.30", "0.00%", "reference"],
                  ["calibrated", "64.30", "0.00%", "control"],
                  ["EFT-MO aug-QZ", "76.14", "+18.41%", "IP baseline"],
                  ["EFT-TDHF aug-QZ 200", "60.73", "-5.55%", "best prediction"],
              ], 11)
    add_bullets(slide, 7.05, 5.0, 5.65, 1.0, ["TDHF 修正了 HF orbital-difference spectrum 的低频权重问题。"], 15, GREEN)

    # 8 TDHF nstates convergence
    slide = prs.slides.add_slide(blank)
    add_title(slide, "15. TDHF 态数收敛：100 起步，推荐 150–200")
    add_picture(slide, save_tdhf_nstates_chart(), 0.65, 1.25, 6.4, 3.65)
    add_table(slide, 7.25, 1.25, 5.25, 3.55,
              ["basis", "nstates", "alpha0", "C6", "C6 err"],
              [
                  ["aug-TZ", "100", "10.3820", "59.7481", "-7.08%"],
                  ["aug-TZ", "200", "10.4224", "60.5445", "-5.84%"],
                  ["aug-QZ", "100", "10.6078", "59.9556", "-6.76%"],
                  ["aug-QZ", "200", "10.6494", "60.7303", "-5.55%"],
              ], 11)
    add_bullets(slide, 0.9, 5.25, 11.7, 1.0, [
        "TRK sum 不必在有限 nstates 内收敛；vdW 低频响应在 100–200 态已经基本稳定。"
    ], 15)

    # 9 Current benchmark
    slide = prs.slides.add_slide(blank)
    set_bg(slide, RGBColor(248, 250, 253))
    add_title(slide, "16. 当前可冻结的 Ar benchmark")
    add_textbox(slide, 1.0, 1.45, 11.5, 0.65, "TDHF aug-cc-pVQZ, nstates = 200", 26, True, BLUE, PP_ALIGN.CENTER)
    add_table(slide, 2.2, 2.45, 8.8, 2.2,
              ["quantity", "value", "reference", "error"],
              [
                  ["α0", "10.64941140", "11.10", "-4.06%"],
                  ["C6", "60.73027908", "64.30", "-5.55%"],
              ], 18)
    add_textbox(slide, 1.0, 5.35, 11.5, 0.55, "这是当前项目第一个可信里程碑：不靠参考反推，直接由 TDHF oscillator spectrum 得到 C6。", 18, True, GREEN, PP_ALIGN.CENTER)

    # 10 Ar2 tail
    slide = prs.slides.add_slide(blank)
    add_title(slide, "17. Ar₂ 长程 tail：E(R) = -C6/R⁶")
    add_picture(slide, save_tail_chart(), 0.7, 1.25, 6.2, 3.65)
    add_table(slide, 7.25, 1.25, 5.1, 3.2,
              ["R (Bohr)", "E_ref (meV)", "E_TDHF", "E_MO"],
              [
                  ["8", "-6.6745", "-6.3040", "-7.9033"],
                  ["10", "-1.7497", "-1.6526", "-2.0718"],
                  ["15", "-0.1536", "-0.1451", "-0.1819"],
                  ["30", "-0.00240", "-0.00227", "-0.00284"],
              ], 11)
    add_bullets(slide, 7.15, 4.85, 5.5, 1.2, [
        "TDHF tail：全程弱 5.55%。",
        "MO tail：全程强 18.41%。"
    ], 15)

    # 11 Relation to EFT paper
    slide = prs.slides.add_slide(blank)
    add_title(slide, "18. 与 EFT-KS 论文的关系")
    add_bullets(slide, 0.75, 1.25, 11.8, 4.8, [
        "论文已经给出 core dynamical Wilson coefficient：dynamic pole + form factor。",
        "本工作把同一 core 动力学项从 band narrowing 推广到 density response / polarizability。",
        "当前 Ar benchmark 暂用 all-electron TDHF total response 验证 α(iξ)→C6 后端。",
        "下一步再回到 frozen-core decomposition：α_total = α_valence + α_core^EFT。"
    ], 19)
    add_code(slide, 1.15, 5.6, 11.0, 0.75, "Eq.(6)/(7) → χ_c(q,q';iξ) → α(iξ) → C6 → E(R)", 20)

    # 12 Code artifacts
    slide = prs.slides.add_slide(blank)
    add_title(slide, "19. 代码与结果文件")
    add_table(slide, 0.55, 1.15, 12.35, 4.65,
              ["类别", "文件"],
              [
                  ["核心后端", "eft_alpha.py, run_alpha_table.py, run_c6_table.py"],
                  ["TDHF 输出", "pyscf_export_ar_tdhf_oscillators.py, run_tdhf_atom.py"],
                  ["收敛测试", "run_tdhf_nstates_convergence_ar.py, run_tdhf_basis_convergence_ar.py"],
                  ["Ar benchmark", "results/ar/ar_tdhf_nstates_convergence.csv, summary.md"],
                  ["Ar2 tail", "run_ar2_tail.py, results/ar/ar2_tail_comparison.csv"],
              ], 12)
    add_bullets(slide, 0.85, 6.0, 11.8, 0.7, ["GitHub 已保存当前 Ar TDHF benchmark 里程碑。"], 15, GREEN)

    # 13 Next steps
    slide = prs.slides.add_slide(blank)
    add_title(slide, "20. 下一步计划")
    add_bullets(slide, 0.85, 1.25, 11.7, 5.2, [
        "1. 用 run_tdhf_atom.py 扩展到 Ne，再到 Kr，先验证闭壳层 noble gas。",
        "2. 建立 Ne/Ar/Kr 的 α0、C6、long-range tail benchmark 表。",
        "3. 若 TDHF noble gas 误差稳定在 5–10%，再进入 frozen-core EFT decomposition。",
        "4. 最终目标：比较 pseudopotential valence-only 漏掉的 core response，以及 EFT core channel 补偿量。",
        "5. D4/MBD 对比先从 long-range C6 开始，不急着做完整 binding curve。"
    ], 18)

    # 14 final
    slide = prs.slides.add_slide(blank)
    set_bg(slide, BLUE)
    add_textbox(slide, 0.9, 1.15, 11.6, 0.9, "结论", 32, True, RGBColor(255, 255, 255), PP_ALIGN.CENTER)
    add_bullets(slide, 1.45, 2.25, 10.8, 3.2, [
        "EFT core dynamics 可以自然连接到 vdW polarizability。",
        "Radial shell-average 路线暴露出 TRK 过计数，暂作诊断。",
        "3D MO 路线给出稳定 baseline，TDHF/RPA 进一步修正频谱形状。",
        "Ar TDHF aug-cc-pVQZ nstates=200：C6 error = -5.55%，形成第一个可信 benchmark。"
    ], 20, RGBColor(255, 255, 255))

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    make_ppt()
